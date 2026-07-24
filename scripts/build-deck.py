#!/usr/bin/env python3
"""Milestone PayFi — presentation deck generator (v2, agentic layer).

Design system pulled from the product (app/globals.css):
  ink #101312, canvas #F5F7F4, accent #245CFF, muted #66727A.
Style: clean/modern — one idea per slide, huge type, keywords only.

Usage:  python3 scripts/build-deck.py   (requires: pip install python-pptx)
Output: milestone-payfi-deck.pptx in the repo root.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- palette ----------
INK = RGBColor(0x10, 0x13, 0x12)
INK_2 = RGBColor(0x1B, 0x20, 0x1E)        # raised surface on dark
LINE_D = RGBColor(0x2C, 0x34, 0x30)       # hairline on dark
PAPER = RGBColor(0xF5, 0xF7, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE_L = RGBColor(0xD9, 0xE1, 0xDD)       # hairline on light
ACCENT = RGBColor(0x24, 0x5C, 0xFF)
ACCENT_D = RGBColor(0x8F, 0xAE, 0xFF)     # accent readable on ink
MUTED_L = RGBColor(0x66, 0x72, 0x7A)      # muted on light
MUTED_D = RGBColor(0xAE, 0xB9, 0xC2)      # muted on dark
MINT = RGBColor(0x78, 0xE0, 0xC4)         # proof green on dark
GREEN = RGBColor(0x08, 0x76, 0x5E)
ON_ACCENT_SUB = RGBColor(0xD5, 0xDE, 0xFF)

SANS = "Arial"
MONO = "Courier New"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
MX = Inches(0.92)          # side margin
CW = EMU_W - 2 * MX        # content width

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide(bg):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def txt(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dicts {runs:[(text, size, color, bold, font)], align, before, after, line}"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, p in enumerate(paras):
        pr = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pr.alignment = p.get("align", PP_ALIGN.LEFT)
        if p.get("before") is not None:
            pr.space_before = Pt(p["before"])
        if p.get("after") is not None:
            pr.space_after = Pt(p["after"])
        pr.line_spacing = p.get("line", 1.0)
        for t, size, color, bold, font in p["runs"]:
            r = pr.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
    return tb


def rect(s, x, y, w, h, fill, line=None, round_=False, radius=0.12, line_w=1.0):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    if round_:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def kicker(s, dark, num, label):
    c_mut = MUTED_D if dark else MUTED_L
    y = Inches(0.62)
    rect(s, MX, y + Inches(0.03), Inches(0.34), Inches(0.16), ACCENT)
    txt(s, MX + Inches(0.52), y - Inches(0.06), Inches(9), Inches(0.4), [
        {"runs": [(num, 13, ACCENT_D if dark else ACCENT, True, MONO),
                  ("   " + label, 13, c_mut, True, MONO)]}
    ])


def footer(s, dark, n):
    c = MUTED_D if dark else MUTED_L
    rect(s, MX, Inches(7.02), CW, Pt(0.9), LINE_D if dark else LINE_L)
    txt(s, MX, Inches(7.12), Inches(6), Inches(0.3),
        [{"runs": [("MILESTONE PAYFI", 9, c, True, MONO)]}])
    txt(s, EMU_W - MX - Inches(6), Inches(7.12), Inches(6), Inches(0.3),
        [{"runs": [("BUILT ON ARC   ", 9, c, True, MONO), (f"{n:02d}", 9, c, True, MONO)],
          "align": PP_ALIGN.RIGHT}])


def chip(s, x, y, w, h, dark, title, sub=None, tsize=17, ssize=11.5):
    fill = INK_2 if dark else WHITE
    ln = LINE_D if dark else LINE_L
    rect(s, x, y, w, h, fill, ln, round_=True, radius=0.10)
    tc = PAPER if dark else INK
    sc = MUTED_D if dark else MUTED_L
    paras = [{"runs": [(title, tsize, tc, True, SANS)], "line": 1.05}]
    if sub:
        paras.append({"runs": [(sub, ssize, sc, False, SANS)], "before": 6, "line": 1.15})
    txt(s, x + Inches(0.26), y, w - Inches(0.52), h, paras, anchor=MSO_ANCHOR.MIDDLE)


def h1(s, dark, parts, y=Inches(1.28), size=58):
    c = PAPER if dark else INK
    runs = [(t, size, col or c, True, SANS) for t, col in parts]
    txt(s, MX, y, CW, Inches(2.2), [{"runs": runs, "line": 1.02}])


CW3 = Inches(3.7)
GAP3 = Inches(0.24)

# =====================================================================
# 01 — COVER
# =====================================================================
s = slide(INK)
rect(s, 0, Inches(7.28), EMU_W, Inches(0.22), ACCENT)
txt(s, MX, Inches(0.62), Inches(10), Inches(0.4), [
    {"runs": [("ENCODE CLUB × ARC", 13, MUTED_D, True, MONO)]}])
txt(s, MX, Inches(2.05), CW, Inches(2.6), [
    {"runs": [("MILESTONE", 88, PAPER, True, SANS)], "line": 0.98},
    {"runs": [("PAYFI", 88, ACCENT_D, True, SANS)], "line": 0.98, "before": 2},
])
txt(s, MX, Inches(4.82), CW, Inches(0.7), [
    {"runs": [("Approved work  ", 26, PAPER, False, SANS),
              ("→", 26, ACCENT_D, True, SANS),
              ("  instant working capital.", 26, PAPER, False, SANS)]}])
txt(s, MX, Inches(5.66), CW, Inches(0.4), [
    {"runs": [("USDC ESCROW × RECEIVABLE POOL × AGENTS   ·   LIVE ON ARC TESTNET", 12, MUTED_D, True, MONO)]}])

# =====================================================================
# 02 — PROBLEM
# =====================================================================
s = slide(PAPER)
kicker(s, False, "01", "THE PROBLEM")
h1(s, False, [("Approved ", None), ("≠", ACCENT), (" paid.", None)])
txt(s, MX, Inches(2.55), CW, Inches(1.0), [
    {"runs": [("Work accepted on Monday.", 22, INK, False, SANS)], "line": 1.15},
    {"runs": [("Cash in Net-30. Net-60. ", 22, MUTED_L, False, SANS),
              ("Net-never.", 22, INK, True, SANS)], "line": 1.15, "before": 4},
])
y = Inches(4.15)
for i, (t, sub) in enumerate([
    ("WAITING", "delivery done, cash locked"),
    ("NO COLLATERAL", "an invoice banks won't touch"),
    ("NO CREDIT FILE", "thin-file by design"),
]):
    chip(s, MX + i * (CW3 + GAP3), y, CW3, Inches(1.45), False, t, sub, tsize=20)
footer(s, False, 2)

# =====================================================================
# 03 — STATUS QUO
# =====================================================================
s = slide(PAPER)
kicker(s, False, "01", "THE PROBLEM")
h1(s, False, [("Today's options.", None)])
txt(s, MX, Inches(2.5), CW, Inches(0.5), [
    {"runs": [("Every route costs time, margin, or dignity.", 20, MUTED_L, False, SANS)]}])
y = Inches(3.6)
for i, (t, l1, l2) in enumerate([
    ("WAIT IT OUT", "free · slow", "rent is due now"),
    ("FACTORING", "paperwork · minimums", "built for big balance sheets"),
    ("PAYDAY APPS", "wages, not invoices", "wrong asset entirely"),
]):
    x = MX + i * (CW3 + GAP3)
    rect(s, x, y, CW3, Inches(2.15), WHITE, LINE_L, round_=True, radius=0.08)
    txt(s, x + Inches(0.28), y + Inches(0.3), CW3 - Inches(0.56), Inches(1.6), [
        {"runs": [(t, 19, INK, True, SANS)]},
        {"runs": [(l1, 13.5, ACCENT, True, MONO)], "before": 10},
        {"runs": [(l2, 13.5, MUTED_L, False, SANS)], "before": 5},
    ])
footer(s, False, 3)

# =====================================================================
# 04 — INSIGHT
# =====================================================================
s = slide(INK)
kicker(s, True, "02", "THE INSIGHT")
h1(s, True, [("Approval flips the risk.", None)])
y = Inches(3.0)
half = Inches(5.62)
rect(s, MX, y, half, Inches(1.9), INK_2, LINE_D, round_=True, radius=0.08)
txt(s, MX + Inches(0.3), y + Inches(0.28), half - Inches(0.6), Inches(1.4), [
    {"runs": [("BEFORE APPROVAL", 12, MUTED_D, True, MONO)]},
    {"runs": [("delivery risk", 20, PAPER, True, SANS)], "before": 8},
    {"runs": [("escrow protects the client", 13.5, MUTED_D, False, SANS)], "before": 5},
])
x2 = MX + half + Inches(0.3)
rect(s, x2, y, half, Inches(1.9), INK_2, ACCENT_D, round_=True, radius=0.08, line_w=1.4)
txt(s, x2 + Inches(0.3), y + Inches(0.28), half - Inches(0.6), Inches(1.4), [
    {"runs": [("AFTER APPROVAL", 12, ACCENT_D, True, MONO)]},
    {"runs": [("funded USDC claim", 20, PAPER, True, SANS)], "before": 8},
    {"runs": [("financeable · programmable", 13.5, MUTED_D, False, SANS)], "before": 5},
])
txt(s, MX, Inches(5.35), CW, Inches(0.7), [
    {"runs": [("Approved work = a programmable receivable.", 26, ACCENT_D, True, SANS)]}])
footer(s, True, 4)

# =====================================================================
# 05 — PRODUCT
# =====================================================================
s = slide(PAPER)
kicker(s, False, "03", "THE PRODUCT")
h1(s, False, [("One shared task room.", None)])
txt(s, MX, Inches(2.5), CW, Inches(0.5), [
    {"runs": [("Not a dashboard. A room both sides already understand.", 20, MUTED_L, False, SANS)]}])
y = Inches(3.5)
for i, (t, l1, l2) in enumerate([
    ("CLIENT", "funds · approves", "protection until work is accepted"),
    ("FREELANCER", "submits · draws", "cash at approval, not Net-30"),
    ("LIQUIDITY", "advances · earns", "short-duration, escrow-backed"),
]):
    x = MX + i * (CW3 + GAP3)
    rect(s, x, y, CW3, Inches(1.95), WHITE, LINE_L, round_=True, radius=0.08)
    rect(s, x + Inches(0.28), y + Inches(0.32), Inches(0.3), Inches(0.14), ACCENT)
    txt(s, x + Inches(0.28), y + Inches(0.58), CW3 - Inches(0.56), Inches(1.3), [
        {"runs": [(t, 19, INK, True, SANS)]},
        {"runs": [(l1, 13.5, ACCENT, True, MONO)], "before": 9},
        {"runs": [(l2, 13.5, MUTED_L, False, SANS)], "before": 5},
    ])
txt(s, MX, Inches(5.75), CW, Inches(0.4), [
    {"runs": [("CCTP APP KITS — FUND FROM ANY CHAIN   ·   AGENTIC LAYER BUILT IN", 12, MUTED_L, True, MONO)],
     "align": PP_ALIGN.CENTER}])
footer(s, False, 5)

# =====================================================================
# 06 — FLOW
# =====================================================================
s = slide(INK)
kicker(s, True, "04", "HOW IT WORKS")
h1(s, True, [("Five steps. One room.", None)])
steps = [
    ("FUND", "client locks USDC"),
    ("SUBMIT", "evidence in room"),
    ("APPROVE", "receivable born"),
    ("PAY NOW", "agent-scored advance"),
    ("SETTLE", "agent repays pool"),
]
y = Inches(3.35)
n = len(steps)
aw = Inches(0.42)
cw5 = (CW - aw * (n - 1)) / n
for i, (t, sub) in enumerate(steps):
    x = MX + i * (cw5 + aw)
    hot = (i == 3)
    rect(s, x, y, cw5, Inches(1.8), ACCENT if hot else INK_2,
         None if hot else LINE_D, round_=True, radius=0.10)
    txt(s, x + Inches(0.2), y + Inches(0.26), cw5 - Inches(0.4), Inches(1.3), [
        {"runs": [(f"0{i+1}", 11, WHITE if hot else MUTED_D, True, MONO)]},
        {"runs": [(t, 17, WHITE if hot else PAPER, True, SANS)], "before": 7, "line": 1.0},
        {"runs": [(sub, 11.5, ON_ACCENT_SUB if hot else MUTED_D, False, SANS)],
         "before": 6, "line": 1.1},
    ])
    if i < n - 1:
        txt(s, x + cw5 - Inches(0.02), y + Inches(0.62), aw + Inches(0.06), Inches(0.5), [
            {"runs": [("→", 20, ACCENT_D, True, SANS)], "align": PP_ALIGN.CENTER}])
txt(s, MX, Inches(5.6), CW, Inches(0.5), [
    {"runs": [("STATE LIVES ON ARC   ·   AGENTS RUN RISK + SETTLEMENT", 12, MUTED_D, True, MONO)],
     "align": PP_ALIGN.CENTER}])
footer(s, True, 6)

# =====================================================================
# 07 — PAYFI MOMENT
# =====================================================================
s = slide(INK)
kicker(s, True, "04", "HOW IT WORKS")
h1(s, True, [("Approval is the ", None), ("PayFi moment.", ACCENT_D)], size=48)
y = Inches(3.2)
wA, wAr, wB = Inches(4.6), Inches(1.4), Inches(4.6)
xA = MX + Inches(0.6)
rect(s, xA, y, wA, Inches(1.7), INK_2, LINE_D, round_=True, radius=0.08)
txt(s, xA, y, wA, Inches(1.7), [
    {"runs": [("ESCROW", 26, MUTED_D, True, SANS)], "align": PP_ALIGN.CENTER},
    {"runs": [("locked capital", 13, MUTED_D, False, SANS)], "align": PP_ALIGN.CENTER, "before": 6},
], anchor=MSO_ANCHOR.MIDDLE)
txt(s, xA + wA, y + Inches(0.5), wAr, Inches(0.7), [
    {"runs": [("→", 34, ACCENT_D, True, SANS)], "align": PP_ALIGN.CENTER}])
xB = xA + wA + wAr
rect(s, xB, y, wB, Inches(1.7), ACCENT, None, round_=True, radius=0.08)
txt(s, xB, y, wB, Inches(1.7), [
    {"runs": [("RECEIVABLE", 26, WHITE, True, SANS)], "align": PP_ALIGN.CENTER},
    {"runs": [("instantly financeable", 13, ON_ACCENT_SUB, False, SANS)],
     "align": PP_ALIGN.CENTER, "before": 6},
], anchor=MSO_ANCHOR.MIDDLE)
txt(s, MX, Inches(5.45), CW, Inches(0.5), [
    {"runs": [("SAME FUNDS   ·   NEW ASSET   ·   ZERO NEW RISK TO THE CLIENT", 12, MUTED_D, True, MONO)],
     "align": PP_ALIGN.CENTER}])
footer(s, True, 7)

# =====================================================================
# 08 — THREE-SIDED WIN
# =====================================================================
s = slide(PAPER)
kicker(s, False, "05", "WHY IT WORKS")
h1(s, False, [("Three incentives. One pool.", None)])
y = Inches(3.35)
for i, (t, l1, l2) in enumerate([
    ("CLIENT", "pay only for approved work", "escrow protection, unchanged"),
    ("FREELANCER", "cash at approval", "not Net-30 · not a loan shark"),
    ("LP", "short-duration yield", "repaid from funded escrow"),
]):
    x = MX + i * (CW3 + GAP3)
    rect(s, x, y, CW3, Inches(2.15), WHITE, LINE_L, round_=True, radius=0.08)
    txt(s, x + Inches(0.28), y + Inches(0.3), CW3 - Inches(0.56), Inches(1.6), [
        {"runs": [(t, 19, INK, True, SANS)]},
        {"runs": [(l1, 15, GREEN, True, SANS)], "before": 10},
        {"runs": [(l2, 13.5, MUTED_L, False, SANS)], "before": 5},
    ])
footer(s, False, 8)

# =====================================================================
# 09 — RISK ENGINE
# =====================================================================
s = slide(PAPER)
kicker(s, False, "06", "RISK ENGINE")
h1(s, False, [("Not every invoice gets funded.", None)], size=46)
txt(s, MX, Inches(2.4), CW, Inches(0.5), [
    {"runs": [("An onchain policy gates every advance — before liquidity moves.", 20, MUTED_L, False, SANS)]}])
guards = [
    "onchain risk policy", "same-wallet fraud block", "45-day tenor cap",
    "65% utilization cap", "client + freelancer exposure caps", "time-based discount",
]
y0 = Inches(3.5)
cw2 = Inches(5.72); gap2 = Inches(0.24); hh = Inches(1.0)
for i, g in enumerate(guards):
    x = MX + (i % 2) * (cw2 + gap2)
    y = y0 + (i // 2) * (hh + Inches(0.18))
    rect(s, x, y, cw2, hh, WHITE, LINE_L, round_=True, radius=0.14)
    rect(s, x + Inches(0.26), y + hh / 2 - Inches(0.07), Inches(0.14), Inches(0.14), ACCENT)
    txt(s, x + Inches(0.6), y, cw2 - Inches(0.8), hh,
        [{"runs": [(g, 16, INK, True, SANS)]}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, False, 9)

# =====================================================================
# 10 — AGENTIC LAYER (new)
# =====================================================================
s = slide(INK)
kicker(s, True, "07", "THE AGENTIC LAYER")
h1(s, True, [("No human ", None), ("in the loop.", ACCENT_D)])
txt(s, MX, Inches(2.55), CW, Inches(0.5), [
    {"runs": [("Risk and settlement run themselves — with revocable onchain delegation.", 20, MUTED_D, False, SANS)]}])
y = Inches(3.45)
for i, (t, l1, l2) in enumerate([
    ("UNDERWRITER", "scores · publishes policy", "signals: liquidity · exposure · tenor · fraud"),
    ("SETTLER", "closes · repays pool", "permissionless release at maturity"),
    ("IDENTITY", "ERC-8004 on Arc", "agent 851709 · owner can revoke"),
]):
    x = MX + i * (CW3 + GAP3)
    rect(s, x, y, CW3, Inches(2.1), INK_2, LINE_D, round_=True, radius=0.08)
    rect(s, x + Inches(0.28), y + Inches(0.3), Inches(0.3), Inches(0.14), ACCENT)
    txt(s, x + Inches(0.28), y + Inches(0.56), CW3 - Inches(0.56), Inches(1.4), [
        {"runs": [(t, 17, PAPER, True, SANS)]},
        {"runs": [(l1, 13, ACCENT_D, True, MONO)], "before": 9},
        {"runs": [(l2, 12, MUTED_D, False, SANS)], "before": 5, "line": 1.2},
    ])
footer(s, True, 10)

# =====================================================================
# 11 — WHY ARC
# =====================================================================
s = slide(INK)
kicker(s, True, "08", "WHY ARC")
txt(s, MX, Inches(1.28), CW, Inches(2.2), [
    {"runs": [("Programmable money,", 54, PAPER, True, SANS)], "line": 1.02},
    {"runs": [("real settlement.", 54, ACCENT_D, True, SANS)], "line": 1.02, "before": 2},
])
y = Inches(3.6)
feats = [
    ("USDC-NATIVE", "the asset is the rail — gas too"),
    ("EVM-COMPATIBLE", "standard tooling, no rewrite"),
    ("DETERMINISTIC", "sub-second finality for agents"),
    ("APP KITS", "CCTP: fund from any chain"),
]
cw4 = (CW - 3 * Inches(0.24)) / 4
for i, (t, sub) in enumerate(feats):
    x = MX + i * (cw4 + Inches(0.24))
    rect(s, x, y, cw4, Inches(1.9), INK_2, LINE_D, round_=True, radius=0.10)
    txt(s, x + Inches(0.24), y + Inches(0.28), cw4 - Inches(0.48), Inches(1.4), [
        {"runs": [(t, 15.5, ACCENT_D, True, MONO)], "line": 1.05},
        {"runs": [(sub, 12.5, MUTED_D, False, SANS)], "before": 9, "line": 1.15},
    ])
footer(s, True, 11)

# =====================================================================
# 12 — LIVE PROOF (v3)
# =====================================================================
s = slide(INK)
kicker(s, True, "09", "LIVE PROOF")
h1(s, True, [("Deployed. Verified. ", None), ("Agent-run.", None)], size=52)
y = Inches(2.85)
rect(s, MX, y, CW, Inches(3.0), INK_2, LINE_D, round_=True, radius=0.05)
pad = Inches(0.4)
mono_rows = [
    ("CHAIN", "Arc Testnet · 5042002", MUTED_D),
    ("ESCROW", "0x605d5f089a27c6a4f7b1271bdc27d03e4336e314", PAPER),
    ("POOL", "0xc1fdb1507f489b5d426f4da398fd4da9d12e108f", PAPER),
    ("AGENT", "underwriter 0x3C06…093E2 · ERC-8004 id 851709", PAPER),
    ("VERIFY", "fund → submit → approve → agent policy → advance → agent settle", MUTED_D),
    ("RESULT", "8 / 8 onchain steps PASS   ·   13 / 13 contract tests PASS", MINT),
]
ry = y + Inches(0.28)
for label, val, vc in mono_rows:
    txt(s, MX + pad, ry, Inches(1.5), Inches(0.35),
        [{"runs": [(label, 12, ACCENT_D, True, MONO)]}])
    txt(s, MX + pad + Inches(1.55), ry, CW - pad * 2 - Inches(1.55), Inches(0.35),
        [{"runs": [(val, 12.5, vc, False, MONO)]}])
    ry += Inches(0.46)
txt(s, MX, y + Inches(3.16), CW, Inches(0.4), [
    {"runs": [("testnet.arcscan.app   ·   pnpm verify:onchain   ·   pnpm contracts:test", 12, MUTED_D, True, MONO)],
     "align": PP_ALIGN.CENTER}])
footer(s, True, 12)

# =====================================================================
# 13 — WHAT'S NEXT
# =====================================================================
s = slide(PAPER)
kicker(s, False, "10", "WHAT'S NEXT")
h1(s, False, [("From demo to desk.", None)])
nxt = [
    ("DURABLE CORE", "Postgres store · Vercel"),
    ("COMPLIANCE", "KYB · KYC · invoices"),
    ("CREDIT", "richer agent signals"),
    ("LIQUIDITY", "term sheets"),
    ("MAINNET", "Arc production"),
]
y = Inches(3.4)
cw5b = (CW - 4 * Inches(0.22)) / 5
for i, (t, sub) in enumerate(nxt):
    x = MX + i * (cw5b + Inches(0.22))
    rect(s, x, y, cw5b, Inches(1.75), WHITE, LINE_L, round_=True, radius=0.10)
    txt(s, x + Inches(0.22), y + Inches(0.26), cw5b - Inches(0.44), Inches(1.3), [
        {"runs": [(t, 14.5, INK, True, SANS)], "line": 1.05},
        {"runs": [(sub, 12, MUTED_L, False, SANS)], "before": 8, "line": 1.15},
    ])
footer(s, False, 13)

# =====================================================================
# 14 — CLOSE
# =====================================================================
s = slide(INK)
rect(s, 0, Inches(7.28), EMU_W, Inches(0.22), ACCENT)
txt(s, MX, Inches(0.62), Inches(10), Inches(0.4), [
    {"runs": [("MILESTONE PAYFI", 13, MUTED_D, True, MONO)]}])
txt(s, MX, Inches(2.5), CW, Inches(2.2), [
    {"runs": [("Approved work should pay", 54, PAPER, True, SANS)], "line": 1.04},
    {"runs": [("like it's settled.", 54, ACCENT_D, True, SANS)], "line": 1.04, "before": 2},
])
txt(s, MX, Inches(5.15), CW, Inches(0.5), [
    {"runs": [("The wedge: approved work becomes instant working capital — run by agents.", 19, MUTED_D, False, SANS)]}])
txt(s, MX, Inches(5.85), CW, Inches(0.4), [
    {"runs": [("LIVE ON ARC TESTNET   ·   ESCROW × POOL × AGENTS   ·   8/8 ONCHAIN STEPS VERIFIED", 12, MUTED_D, True, MONO)]}])

prs.core_properties.title = "Milestone PayFi — Deck"
prs.core_properties.author = "Milestone PayFi"
OUT = "milestone-payfi-deck.pptx"
prs.save(OUT)
print(f"saved {OUT} ({len(prs.slides._sldIdLst)} slides)")
